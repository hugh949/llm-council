"""Document parsing utilities for extracting text from various file formats."""

import io
from typing import Optional, Dict, Any
import httpx


async def parse_pdf(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from PDF file."""
    try:
        import pypdf
        
        pdf_file = io.BytesIO(file_content)
        pdf_reader = pypdf.PdfReader(pdf_file)
        
        text_parts = []
        for page_num, page in enumerate(pdf_reader.pages, 1):
            try:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{text}")
            except Exception as e:
                text_parts.append(f"--- Page {page_num} (Error extracting text: {str(e)}) ---")
        
        return {
            "type": "pdf",
            "name": filename,
            "content": "\n\n".join(text_parts),
            "page_count": len(pdf_reader.pages)
        }
    except ImportError:
        return {"type": "pdf", "name": filename, "content": "[PDF parsing not available]", "error": "pypdf not installed"}
    except Exception as e:
        return {"type": "pdf", "name": filename, "content": f"[Error parsing PDF: {str(e)}]", "error": str(e)}


async def parse_docx(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from Word document."""
    try:
        from docx import Document
        
        doc_file = io.BytesIO(file_content)
        doc = Document(doc_file)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n".join(table_text))
        
        return {
            "type": "docx",
            "name": filename,
            "content": "\n\n".join(text_parts)
        }
    except ImportError:
        return {"type": "docx", "name": filename, "content": "[Word document parsing not available]", "error": "python-docx not installed"}
    except Exception as e:
        return {"type": "docx", "name": filename, "content": f"[Error parsing Word document: {str(e)}]", "error": str(e)}


async def parse_xlsx(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from Excel spreadsheet."""
    try:
        from openpyxl import load_workbook
        
        excel_file = io.BytesIO(file_content)
        workbook = load_workbook(excel_file, data_only=True)
        
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_values):
                    sheet_text.append(" | ".join(row_values))
            
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
        
        return {
            "type": "xlsx",
            "name": filename,
            "content": "\n\n".join(text_parts)
        }
    except ImportError:
        return {"type": "xlsx", "name": filename, "content": "[Excel parsing not available]", "error": "openpyxl not installed"}
    except Exception as e:
        return {"type": "xlsx", "name": filename, "content": f"[Error parsing Excel file: {str(e)}]", "error": str(e)}


async def parse_pptx(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text from PowerPoint presentation."""
    try:
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return {
            "type": "pptx",
            "name": filename,
            "content": "\n\n".join(text_parts),
            "slide_count": len(prs.slides)
        }
    except ImportError:
        return {"type": "pptx", "name": filename, "content": "[PowerPoint parsing not available]", "error": "python-pptx not installed"}
    except Exception as e:
        return {"type": "pptx", "name": filename, "content": f"[Error parsing PowerPoint: {str(e)}]", "error": str(e)}


async def parse_text_file(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Parse plain text file."""
    try:
        # Try UTF-8 first
        try:
            content = file_content.decode('utf-8')
        except UnicodeDecodeError:
            # Fall back to latin-1
            content = file_content.decode('latin-1')
        
        return {
            "type": "text",
            "name": filename,
            "content": content
        }
    except Exception as e:
        return {"type": "text", "name": filename, "content": f"[Error reading text file: {str(e)}]", "error": str(e)}


async def parse_file(file_content: bytes, filename: str) -> Dict[str, Any]:
    """
    Parse a file based on its extension.
    
    Args:
        file_content: Binary file content
        filename: Original filename
        
    Returns:
        Dict with type, name, content, and optional metadata
    """
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf'):
        return await parse_pdf(file_content, filename)
    elif filename_lower.endswith(('.doc', '.docx')):
        return await parse_docx(file_content, filename)
    elif filename_lower.endswith(('.xls', '.xlsx')):
        return await parse_xlsx(file_content, filename)
    elif filename_lower.endswith(('.ppt', '.pptx')):
        return await parse_pptx(file_content, filename)
    elif filename_lower.endswith(('.txt', '.md', '.csv')):
        return await parse_text_file(file_content, filename)
    else:
        # Try to parse as text as fallback
        return await parse_text_file(file_content, filename)


async def fetch_url_content(url: str) -> Dict[str, Any]:
    """
    Fetch and extract text content from a URL.
    
    Args:
        url: URL to fetch
        
    Returns:
        Dict with type, name (url), content
    """
    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            response = await client.get(url, headers={
                "User-Agent": "Mozilla/5.0 (compatible; LLM-Council/1.0)"
            })
            response.raise_for_status()
            
            content_type = response.headers.get("content-type", "").lower()
            
            if "text/html" in content_type:
                # Parse HTML
                try:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(response.text, 'lxml')
                    
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    # Get text
                    text = soup.get_text(separator="\n", strip=True)
                    
                    # Clean up extra whitespace
                    lines = [line.strip() for line in text.split("\n") if line.strip()]
                    text = "\n".join(lines)
                    
                    return {
                        "type": "url",
                        "name": url,
                        "content": text[:50000],  # Limit to 50k characters
                        "original_url": url
                    }
                except ImportError:
                    return {
                        "type": "url",
                        "name": url,
                        "content": response.text[:50000],
                        "original_url": url
                    }
            else:
                # For non-HTML, return as text
                return {
                    "type": "url",
                    "name": url,
                    "content": response.text[:50000],
                    "original_url": url
                }
    except Exception as e:
        return {
            "type": "url",
            "name": url,
            "content": f"[Error fetching URL: {str(e)}]",
            "original_url": url,
            "error": str(e)
        }


